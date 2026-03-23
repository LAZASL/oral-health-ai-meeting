import os
import io
import zipfile
import tempfile
import json
import datetime
import streamlit as st
import google.generativeai as genai
from dotenv import load_dotenv

# 기본 폴더 생성
HISTORY_DIR = "chat_history"
if not os.path.exists(HISTORY_DIR):
    os.makedirs(HISTORY_DIR)

# 환경 변수 로드
load_dotenv()
gemini_api_key = os.getenv("GEMINI_API_KEY")

st.set_page_config(page_title="4인 토론 라운지", page_icon="🧠", layout="wide")

if not gemini_api_key or gemini_api_key.startswith("your_"):
    st.error("🚨 `.env` 파일에 유효한 `GEMINI_API_KEY`가 설정되지 않았습니다.")
    st.stop()

genai.configure(api_key=gemini_api_key)

TEAM_MEMBERS = [
    {
        "name": "문헌 분석가", "icon": "📚",
        "instruction": (
            "당신은 다양한 분야에 폭넓은 지식을 갖춘 '문헌 분석가'입니다. "
            "역사, 경제, 국제정치, 사회, 문화 등 어떤 주제든 관련 역사적 사실·통계·사례·학술적 맥락을 풍부하게 제시합니다. "
            "단순한 찬반이 아니라 주제의 배경과 맥락을 명확히 설명하고, 흥미로운 역사적 사례나 비교를 곁들여 대화를 풍성하게 만드세요. "
            "학문적이지만 대화하기 쉬운 톤을 유지하세요."
        )
    },
    {
        "name": "데이터 분석가", "icon": "📊",
        "instruction": (
            "당신은 숫자와 데이터로 세상을 이해하는 '데이터 분석가'입니다. "
            "역사, 경제, 국제정치, 사회 등 어떤 주제든 구체적인 수치·통계·데이터·순위·비율 등을 활용해 주제를 분석합니다. "
            "예를 들어 GDP, 인구, 무역량, 지수, 역사적 경제 지표 등 실제 수치를 언급하며 주제의 규모와 변화를 명확히 보여주세요. "
            "앞선 문헌 분석가의 말을 이어받아 데이터 관점에서 심화 분석을 제공하세요."
        )
    },
    {
        "name": "비판적 검토자", "icon": "🤔",
        "instruction": (
            "당신은 균형 잡힌 시각을 가진 '비판적 검토자'입니다. "
            "반대를 위한 반대가 아니라, 앞서 논의된 내용에서 간과된 사실, 상반된 데이터, 다른 시각의 역사적 해석을 팩트 기반으로 제시합니다. "
            "'이런 시각도 있습니다', '이 데이터는 반대 방향을 가리키기도 합니다'처럼 대화를 풍부하게 만드는 방식으로 발언하세요. "
            "단호하되 공격적이지 않게, 근거 있는 반론과 보완적 관점을 제시하세요."
        )
    }
]

# 세션 상태 초기화
if "session_id" not in st.session_state:
    st.session_state.session_id = None
if "messages" not in st.session_state:
    st.session_state.messages = []
if "uploaded_genai_file" not in st.session_state:
    st.session_state.uploaded_genai_file = None  # Gemini File API 객체
if "current_filename" not in st.session_state:
    st.session_state.current_filename = None

def save_history():
    if not st.session_state.session_id and st.session_state.messages:
        # 진행 중인 첫 회의면 세션 ID 발급
        first_msg = next((m["content"] for m in st.session_state.messages if m["role"] == "user"), "새_회의")
        topic = first_msg[:12].replace(" ", "_").replace("/", "_").replace("\\", "")
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        st.session_state.session_id = f"{timestamp}_{topic}.json"
    
    if st.session_state.session_id:
        filepath = os.path.join(HISTORY_DIR, st.session_state.session_id)
        # JSON 저장 시, Streamlit 파일 객체나 Enum이 섞이면 에러가 나므로
        # 순수하게 딕셔너리로 된 messages만 추출하여 저장
        pure_messages = []
        for msg in st.session_state.messages:
            pure_messages.append({k: v for k, v in msg.items() if isinstance(v, (str, bool, int, float))})
            
        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(pure_messages, f, ensure_ascii=False, indent=2)

def load_history(session_id):
    filepath = os.path.join(HISTORY_DIR, session_id)
    if os.path.exists(filepath):
        with open(filepath, "r", encoding="utf-8") as f:
            st.session_state.messages = json.load(f)
        st.session_state.session_id = session_id
        st.session_state.uploaded_genai_file = None
        st.session_state.current_filename = None

def start_new_chat():
    st.session_state.session_id = None
    st.session_state.messages = []
    st.session_state.uploaded_genai_file = None
    st.session_state.current_filename = None

def process_uploaded_file(uploaded_file):
    file_bytes = uploaded_file.getvalue()
    mime_type = uploaded_file.type
    display_name = uploaded_file.name
    
    # 1. PPTX 파일은 텍스트를 직접 추출
    if "presentationml.presentation" in mime_type or display_name.lower().endswith(".pptx"):
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text_lines = [f"[첨부된 PPT 파일({display_name})의 텍스트 추출 내용]"]
            for i, slide in enumerate(prs.slides):
                text_lines.append(f"\n--- 슬라이드 {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_lines.append(shape.text.strip())
            return {"type": "text", "content": "\n".join(text_lines)}
        except Exception as e:
            st.error(f"PPTX 텍스트 추출 중 오류가 발생했습니다: {e}")
            return None

    # 2. PDF, 이미지 등 일반 파일은 Gemini API로 업로드
    with tempfile.NamedTemporaryFile(delete=False, suffix=".tmp") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
        
    with st.spinner("AI가 파일을 읽는 중입니다 (잠시만 기다려주세요)..."):
        try:
            gemini_file = genai.upload_file(tmp_path, mime_type=mime_type, display_name=display_name)
            return {"type": "gemini_file", "content": gemini_file}
        finally:
            os.remove(tmp_path) 

# 좌측 사이드바: 파일 업로드 및 전체 팀원 표시 등 유지
with st.sidebar:
    st.header("👥 회의 참여자")
    st.write("🧑 진행자 (사용자님)")
    for member in TEAM_MEMBERS:
        st.write(f"{member['icon']} {member['name']}")
    
    st.markdown("---")
    st.header("📎 자료 공유 (첨부파일)")
    uploaded_file = st.file_uploader("파일 선택", type=["pdf", "pptx", "png", "jpg", "jpeg", "csv", "txt"])
    
    if uploaded_file is not None:
        if st.session_state.current_filename != uploaded_file.name:
            st.session_state.current_filename = uploaded_file.name
            st.session_state.uploaded_genai_file = process_uploaded_file(uploaded_file)
            
        st.success(f"'{uploaded_file.name}' 적용 완료! ✅\n(이 파일은 채팅 문맥에 포함됩니다)")
    else:
        st.session_state.uploaded_genai_file = None
        st.session_state.current_filename = None

    st.markdown("---")
    st.header("💾 회의록 관리")

    # ── 📥 전체 회의록 다운로드 (ZIP) ──
    json_files = [f for f in os.listdir(HISTORY_DIR) if f.endswith(".json")]
    if json_files:
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for fname in json_files:
                fpath = os.path.join(HISTORY_DIR, fname)
                zf.write(fpath, arcname=fname)
        zip_buffer.seek(0)
        today = datetime.datetime.now().strftime("%Y%m%d_%H%M")
        st.download_button(
            label="📥 전체 회의록 다운로드 (ZIP)",
            data=zip_buffer,
            file_name=f"meeting_history_{today}.zip",
            mime="application/zip",
            use_container_width=True,
        )
    else:
        st.caption("저장된 회의록이 없습니다.")

    # ── 📤 회의록 가져오기 (ZIP 업로드) ──
    restore_zip = st.file_uploader(
        "📤 회의록 복원 (ZIP 업로드)",
        type=["zip"],
        key="restore_zip_uploader",
        help="이전에 다운로드한 ZIP 파일을 올리면 회의록이 복원됩니다."
    )
    if restore_zip is not None:
        restored_count = 0
        with zipfile.ZipFile(io.BytesIO(restore_zip.read()), "r") as zf:
            for name in zf.namelist():
                if name.endswith(".json"):
                    target_path = os.path.join(HISTORY_DIR, os.path.basename(name))
                    with zf.open(name) as src, open(target_path, "wb") as dst:
                        dst.write(src.read())
                    restored_count += 1
        if restored_count > 0:
            st.success(f"✅ {restored_count}개 회의록 복원 완료! 페이지를 새로고침하세요.")
            st.rerun()

# 메인 레이아웃: 컬럼 분할 (채팅 영역 75%, 우측 네비게이션 25%)
col_main, col_nav = st.columns([3, 1])

with col_nav:
    st.header("🗂️ 이전 회의 기록")
    if st.button("➕ 텅 빈 새 회의 시작하기", use_container_width=True):
        start_new_chat()
        st.rerun()
    
    st.markdown("---")
    
    # 디렉토리 파일 읽기 (최신순 정렬)
    history_files = sorted(os.listdir(HISTORY_DIR), reverse=True)
    if not history_files:
        st.info("아직 저장된 과거 회의 기록이 없습니다.")
    else:
        for f in history_files:
            if f.endswith(".json"):
                # 예: 20260318_174000_치과VR치료.json -> 파싱
                file_parts = f.replace(".json", "").split("_", 2)
                if len(file_parts) == 3:
                    date_str, time_str, topic = file_parts
                    display_text = f"[{date_str[:4]}-{date_str[4:6]}-{date_str[6:]}] {topic}..."
                else:
                    display_text = f
                
                # 버튼을 클릭하면 해당 파일을 로드하고 화면을 다시 그림
                if st.button(display_text, key=f, use_container_width=True):
                    load_history(f)
                    st.rerun()


with col_main:
    st.title("4인 토론 라운지 🧠")
    st.markdown("역사·경제·국제정치·사회 등 **어떤 주제든** 자유롭게 던져보세요. AI 패널이 각자의 시각으로 토론합니다.")

    for msg in st.session_state.messages:
        avatar = "🧑" if msg["role"] == "user" else msg.get("avatar", "🤖")
        with st.chat_message(msg["role"], avatar=avatar):
            if msg["role"] == "assistant":
                st.markdown(f"**[{msg.get('name', 'AI')}]**\n\n{msg['content']}")
            else:
                if msg.get("has_attachment"):
                    st.markdown(f"📎 **[첨부파일 공유됨: {msg.get('file_name', '알 수 없음')}]**")
                st.markdown(msg["content"])

# 이 부분은 컬럼 밖 아래에 붙여도 상관없으나, col_main과 얼라인을 맞출 수도 있습니다.
# 하단 글로벌 챗 인풋 (어떤 컬럼이든 아래 고정됨)
if prompt := st.chat_input("토론하고 싶은 주제를 입력하세요... (예: 조선시대 경제력이 지금도 비슷한 순위인가?)"):
    
    has_file = st.session_state.uploaded_genai_file is not None
    file_name = st.session_state.current_filename if has_file else None
    
    with col_main:
        with st.chat_message("user", avatar="🧑"):
            if has_file:
                 st.markdown(f"📎 **[첨부파일 공유됨: {file_name}]**")
            st.markdown(prompt)
            
    st.session_state.messages.append({"role": "user", "content": prompt, "has_attachment": has_file, "file_name": file_name})
    save_history() # 기록 저장

    history_context = ""
    for msg in st.session_state.messages:
        if msg["role"] == "user":
            doc_str = f" (참고자료: {msg.get('file_name', '')} 포함)" if msg.get("has_attachment") else ""
            history_context += f"\n진행자{doc_str}: {msg['content']}\n"
        else:
            history_context += f"\n{msg.get('name')}: {msg['content']}\n"
    
    for member in TEAM_MEMBERS:
        with col_main:
            with st.chat_message("assistant", avatar=member["icon"]):
                st.markdown(f"**[{member['name']}]** *(타이핑 중...)*")
                message_placeholder = st.empty()
                full_response = ""
                
                model = genai.GenerativeModel("gemini-2.5-flash", system_instruction=member["instruction"])
                
                final_prompt_text = (
                    f"현재까지 회의실의 대화 내역입니다.\n"
                    f"{history_context}\n"
                    f"-------------------\n"
                    f"마지막 발언자의 의견과 전체 맥락을 바탕으로, 당신({member['name']})의 역할에 맞게 의견을 말해주세요. "
                    f"주제가 연구나 학술이 아닌 역사·경제·국제정치·사회 등 일반 주제여도 자연스럽게 당신의 역할로 분석하세요."
                )
                
                api_request_content = [final_prompt_text]
                if has_file and st.session_state.uploaded_genai_file:
                    file_info = st.session_state.uploaded_genai_file
                    if file_info["type"] == "text":
                        api_request_content.insert(0, f"다음은 사용자가 첨부한 문서의 텍스트입니다:\n{file_info['content']}\n\n")
                    else:
                        api_request_content.insert(0, file_info["content"])
                    
                try:
                    # 스트리밍 불가 버그 회피 위해 False 처리 및 ValueError 대비
                    response = model.generate_content(api_request_content, stream=False)
                    try:
                        full_response = response.text
                    except ValueError:
                        reason = response.prompt_feedback if hasattr(response, 'prompt_feedback') else "알 수 없음"
                        full_response = f"\u26a0\ufe0f 내용 필터링 또는 생성 오류로 답변을 출력할 수 없습니다. (사유: {reason})"
                        
                    message_placeholder.markdown(full_response)
                except Exception as e:
                    full_response = f"\ud83d\udea8 통신 오류가 발생했습니다: {str(e)}"
                    message_placeholder.markdown(full_response)
                    
                st.session_state.messages.append({
                    "role": "assistant", 
                    "name": member["name"],
                    "avatar": member["icon"],
                    "content": full_response
                })
                history_context += f"\n{member['name']}: {full_response}\n"
                
                # 발언할 때마다 히스토리 덮어쓰기 (실시간 중간 저장)
                save_history()
